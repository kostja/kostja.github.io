#!/usr/bin/env python3
"""
Post-process the YDB positioning deck.

ydb-positioning.pptx is generated elsewhere by PptxGenJS (creator "Stroppy").
This one-shot patcher takes that raw deck and:
  - fixes typos (slide 1 title, slide 4 roadmap text);
  - translates and tidies the Russian-language copy;
  - cleans up the "Возможные сценарии go-to-market" slide (renumber 1-6,
    fix zebra striping, recolor badges, drop a leftover orphan box) and
    normalises its row spacing;
  - inserts a "Статус тестирования" progress-bar slide before the scenarios;
  - drops the closing "honest positioning" slide;
  - uses "Greengage" (never Greenplum/Cloudberry).

Usage:
    python3 talks/build_ydb.py [SRC.pptx] [OUT.pptx]

Defaults reproduce the shipped deck from its preserved original:
    SRC = ~/ydb-positioning.orig.pptx   (untouched PptxGenJS export)
    OUT = ~/ydb-positioning.pptx
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = sys.argv[1] if len(sys.argv) > 1 else os.path.expanduser('~/ydb-positioning.orig.pptx')
OUT = sys.argv[2] if len(sys.argv) > 2 else os.path.expanduser('~/ydb-positioning.pptx')

prs = Presentation(SRC)

# ---------- helpers ----------
def C(hexs): return RGBColor.from_string(hexs)

def by_name(slide):
    return {sh.name: sh for sh in slide.shapes}

def set_para(p, text):
    """Set paragraph text into its first run, drop extra runs, keep run0 font."""
    if not p.runs:
        r = p.add_run(); r.text = text; return
    p.runs[0].text = text
    for r in list(p.runs)[1:]:
        r._r.getparent().remove(r._r)

def set_single(shape, text):
    set_para(shape.text_frame.paragraphs[0], text)

def fill(shape, hexs):
    shape.fill.solid(); shape.fill.fore_color.rgb = C(hexs)

def font_color(shape, hexs):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = C(hexs)

def move(shape, x=None, y=None, w=None, h=None):
    if x is not None: shape.left = Inches(x)
    if y is not None: shape.top = Inches(y)
    if w is not None: shape.width = Inches(w)
    if h is not None: shape.height = Inches(h)

def no_line_no_shadow(shape):
    shape.line.fill.background()
    try: shape.shadow.inherit = False
    except Exception: pass

def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[index]
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)

# =========================================================
# SLIDE 1 - title typo
# =========================================================
s1 = by_name(prs.slides[0])
set_single(s1['Text 0'], 'Тестирование YDB')

# =========================================================
# SLIDE 2 - renumber, translate, zebra, colors, dejunk
# =========================================================
s2 = by_name(prs.slides[1])

# delete leftover junk shapes (phantom row + orphan "No")
for nm in ('Text 17', 'Text 19', 'Text 20'):
    sp = s2[nm]; sp._element.getparent().remove(sp._element)

# title + headers
set_single(s2['Text 0'], 'Возможные сценарии go-to-market')
set_single(s2['Text 3'], 'Применимость')

# renumber 1..6
for nm, n in (('Text 5','1'),('Text 11','2'),('Text 22','3'),
              ('Text 27','4'),('Text 33','5'),('Text 38','6')):
    set_single(s2[nm], n)

# scenario names (translate the non-Russian ones)
set_single(s2['Text 23'], 'OLAP')
set_single(s2['Text 39'], 'Против PG / Greengage')

# fit badges -> Russian
set_single(s2['Text 8'],  'Нет')
set_single(s2['Text 14'], 'At scale')  # kept in English by request
set_single(s2['Text 25'], 'Средне')
set_single(s2['Text 30'], 'Opportunity')
set_single(s2['Text 36'], 'Нет')
# 'Не применимо' already Russian

# pill colors: Mid-tier green->amber; Overstated yellow->gold + dark text
fill(s2['Shape 24'], 'B45309')           # Средне -> amber (white text OK)
fill(s2['Shape 29'], 'FACC15')           # Opportunity -> gold
font_color(s2['Text 30'], '1E293B')      # dark text on gold

# notes -> Russian
set_single(s2['Text 9'],
    '~10× медленнее нативного; POOL_SIZE=1, bulkSize=20. Только как средство '
    'миграции. Разрыв в экосистеме (инструменты, зрелость планировщика, ORM, '
    'кадры) носит структурный характер.')
set_single(s2['Text 15'],
    'Одноузловой PG выигрывает в 2–3×. YDB выигрывает, когда запись не помещается '
    'на один сервер или нужны субсекундное переключение при отказе / мультирегиональность.')
set_single(s2['Text 26'],
    'Отстаёт от ClickHouse / DuckDB / Snowflake. Q21 — 72 с на SF100. '
    'Резкая деградация на SF100 при том же оборудовании.')
set_single(s2['Text 31'],
    'Строковые и колоночные таблицы в одном кластере, но нет автоматической '
    'репликации строки→колонки. Ручной конвейер = то же, что две отдельные системы.')
set_single(s2['Text 37'],
    'Нет каталога Iceberg/Delta/Hudi. Коннектор Spark экспериментальный, '
    'колоночные таблицы не поддерживаются.')
set_single(s2['Text 42'],
    'Конкурирует с TiDB / Cockroach / Spanner в сегменте распределённого SQL, '
    'а не с PG или Greengage.')

# --- uniform row geometry + correct zebra (gray on rows 1,3,5) ---
PITCH = 0.74
Y0 = 1.65
rowY = [Y0 + PITCH * i for i in range(6)]
# columns per row: (number, scenario, pill_bg, pill_txt, note)
rows = [
    ('Text 5','Text 6','Shape 7','Text 8','Text 9'),
    ('Text 11','Text 12','Shape 13','Text 14','Text 15'),
    ('Text 22','Text 23','Shape 24','Text 25','Text 26'),
    ('Text 27','Text 28','Shape 29','Text 30','Text 31'),
    ('Text 33','Text 34','Shape 35','Text 36','Text 37'),
    ('Text 38','Text 39','Shape 40','Text 41','Text 42'),
]
for i,(num,scen,pbg,ptx,note) in enumerate(rows):
    y = rowY[i]
    move(s2[num],  y=y); move(s2[scen], y=y); move(s2[note], y=y)
    move(s2[pbg],  y=y+0.13); move(s2[ptx], y=y+0.13)

# gray bands -> rows 1,3,5
bands = ('Shape 10','Shape 21','Shape 32')
for band, ri in zip(bands, (1,3,5)):
    move(s2[band], x=0.55, y=rowY[ri]-0.06, w=12.2, h=PITCH)

# =========================================================
# SLIDE 3 - typos + translate fragments
# =========================================================
s3 = by_name(prs.slides[2])
set_single(s3['Text 4'],
    'Сценарий: реализовано гибридное хранение или автоматическая репликация строки→колонки')
set_single(s3['Text 9'],
    'Сценарий: новый оптимизатор для OLAP')
set_single(s3['Text 5'],
    'Канонический HTAP. Привлекательна потенциальной экономией TCO за счёт снижения '
    'затрат на ETL и отдельный OLAP-кластер в сегменте замены Exadata и крупных '
    'аналитических СУБД (SAP HANA). ')
set_single(s3['Text 10'],
    'Возможность унификации аналитических и транзакционных задач. Привлекательна для '
    'клиента возможностью консолидации служб эксплуатации и стека разработки. '
    'Промежуточная веха для HTAP.')

t13 = s3['Text 13'].text_frame.paragraphs
# P0 heading stays
set_para(t13[1],
    'Целью является не миграция существующих клиентов PostgreSQL, а развитие '
    'экосистемы разработчиков и снижение порога входа для пользователей и '
    'администраторов баз данных.')
set_para(t13[2],
    'Фокус на HTAP (а не, например, Lakehouse или конкуренция с PostgreSQL/Greengage)')
set_para(t13[3],
    'Риски и недоверие к новому игроку при наличии существующих внедрённых систем '
    'требуют достаточно весомого контр-аргумента, и им может оказаться возможность '
    'мультитенантности и выполнение задач с различным профилем нагрузки в одном '
    'контуре. При этом сама парадигма HTAP может и не реализоваться, но наличие такой '
    'возможности будет привлекательным оружием во внутрикорпоративной войне CIO и CDO.')

# the longer "Сценарий: ..." card headings wrap to 3 lines; lower the
# card descriptions so they don't collide with the headings
move(s3['Text 5'], y=2.45)
move(s3['Text 10'], y=2.45)

# =========================================================
# SLIDE 4 (honest positioning) - dropped by request
# =========================================================
delete_slide(prs, 3)   # original slide index 3 = "honest positioning"

# =========================================================
# NEW SLIDE (progress) -> insert as slide #2 (before scenarios)
# =========================================================
layout = prs.slide_layouts[0]          # DEFAULT layout, no placeholders
ns = prs.slides.add_slide(layout)

def add_text(slide, x, y, w, h, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = C(color)
    return tb

def add_bar(slide, x, y, w, h, hexs):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    fill(sp, hexs); no_line_no_shadow(sp)
    try: sp.adjustments[0] = 0.5
    except Exception: pass
    return sp

# title
add_text(ns, 0.6, 0.35, 12.0, 0.7, 'Статус тестирования на 01.06.2026', 32, '1E293B', bold=True)

# bars span ~2/3 of the slide width
TRACK_X, TRACK_W, BAR_H = 0.6, 7.0, 0.30
PCT_X, PCT_W = 7.8, 1.7
bars = [
    ('Совместимость с PostgreSQL', 0.60, '8B5A2B', '8B5A2B'),  # brown
    ('TPC-C',                      1.00, '16A34A', '16A34A'),  # green
    ('TPC-C managed',              1.00, '16A34A', '16A34A'),  # green
    ('TPC-H',                      0.35, 'EAB308', 'CA8A04'),  # yellow (dark text)
    ('TPC-H managed',              0.00, 'B91C1C', 'B91C1C'),  # red, not started
]
y, PITCH = 1.5, 0.95
for label, pct, barcol, pctcol in bars:
    add_text(ns, 0.6, y, 9.0, 0.35, label, 16, '1E293B', bold=True)
    add_bar(ns, TRACK_X, y+0.42, TRACK_W, BAR_H, 'E2E8F0')          # track
    if pct > 0:
        add_bar(ns, TRACK_X, y+0.42, TRACK_W*pct, BAR_H, barcol)    # fill
    add_text(ns, PCT_X, y+0.30, PCT_W, 0.5, f'{int(pct*100)} %', 18,
             pctcol, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    y += PITCH

add_text(ns, 0.6, y+0.05, 12.0, 0.6,
         'Lakehouse и HTAP — не применимо, тесты не проводились.',
         14, '64748B', bold=False)

# move the new (last) slide to position index 1 (becomes slide #2)
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[-1])
sldIdLst.insert(1, ids[-1])

prs.save(OUT)
print('saved', OUT, 'with', len(prs.slides._sldIdLst), 'slides')
