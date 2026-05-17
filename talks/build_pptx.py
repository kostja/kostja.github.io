#!/usr/bin/env python3
"""Build a PowerPoint (.pptx) deck from the C++ Russia talk markdown.

For each `## NN. Title` slide block, builds a PPTX slide with:
- title at top
- any referenced images (SVG converted to PNG via cairosvg), stacked
  if more than one, sized to fit the remaining height
- bullets below the image (or filling the body if no image)
- the blockquoted speaker notes attached to the slide's notes pane

Run: python3 talks/build_pptx.py
Output: talks/cpp-russia-lsm-compaction.pptx
"""

import os
import re
import hashlib

import cairosvg
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

REPO = "/home/kostja/work/kostja.github.io"
POST = f"{REPO}/_posts/2026-05-17-cpp-russia-lsm-compaction.markdown"
OUT = f"{REPO}/talks/cpp-russia-lsm-compaction.pptx"
PNG_CACHE = f"{REPO}/talks/_pptx_cache"
os.makedirs(PNG_CACHE, exist_ok=True)

# 16:9 slide, 13.333 × 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.4)
MARGIN_R = Inches(0.4)
MARGIN_TOP = Inches(0.3)

TITLE_H = Inches(0.7)
TITLE_PT = 30
BODY_PT = 20
NOTE_PT = 12

INK = RGBColor(0x2B, 0x13, 0x21)


def svg_to_png(svg_path):
    """Render an SVG to PNG via cairosvg, cache by content hash."""
    with open(svg_path, "rb") as f:
        svg = f.read()
    h = hashlib.sha1(svg).hexdigest()[:12]
    png_path = f"{PNG_CACHE}/{os.path.basename(svg_path).replace('.svg','')}_{h}.png"
    if not os.path.exists(png_path):
        # 200 DPI gives sharp images at typical projection sizes
        cairosvg.svg2png(bytestring=svg, write_to=png_path,
                         output_width=2200)
    return png_path


# ── Markdown parsing ─────────────────────────────────────
def parse_post(path):
    text = open(path).read()
    # Drop YAML frontmatter
    text = re.sub(r"^---\n.*?\n---\n", "", text, count=1, flags=re.S)
    # Drop the <style> block
    text = re.sub(r"<style>.*?</style>\n*", "", text, flags=re.S)

    parts = re.split(r"\n---\n", text)
    slides = []
    for p in parts:
        m = re.search(r"^## (\d+)\.\s*(.+)$", p, re.M)
        if not m:
            # Title slide doesn't use ## NN
            if "# " in p and "## 0. Title" in p:
                slides.append(parse_title_slide(p))
            continue
        num = int(m.group(1))
        title = m.group(2).strip()
        slides.append({
            "num": num,
            "title": title,
            "is_title": False,
            "images": [m.group(1) for m in re.finditer(r"!\[[^\]]*\]\(([^)]+)\)", p)],
            "bullets": parse_bullets(p),
            "notes": parse_notes(p),
        })
    return slides


def parse_title_slide(p):
    # Big title is `# ...`, subtitle is `### ...`
    h1 = re.search(r"^# (.+)$", p, re.M)
    h3 = re.search(r"^### (.+)$", p, re.M)
    # Remaining non-blockquote lines after the two headings are the author
    body_lines = []
    for line in p.split("\n"):
        if not line.strip():
            continue
        if line.startswith("#") or line.startswith(">"):
            continue
        if line.startswith("## "):
            continue
        body_lines.append(line.strip())
    notes = parse_notes(p)
    return {
        "num": 0,
        "title": h1.group(1) if h1 else "",
        "subtitle": h3.group(1) if h3 else "",
        "byline": " ".join(body_lines),
        "is_title": True,
        "images": [],
        "bullets": [],
        "notes": notes,
    }


def parse_bullets(p):
    """Return a list of (depth, text) bullets. Skip the title heading."""
    bullets = []
    for line in p.split("\n"):
        m = re.match(r"^(\s*)- (.+)$", line)
        if m:
            depth = len(m.group(1)) // 2
            text = strip_md_inline(m.group(2))
            bullets.append((depth, text))
        elif bullets and line.startswith("  ") and not line.startswith("  -"):
            # Continuation of previous bullet (e.g. wrapped text)
            depth, text = bullets[-1]
            bullets[-1] = (depth, text + " " + strip_md_inline(line.strip()))
        elif line.startswith("**") and line.endswith("**"):
            # Section header within bullet list (References slide)
            bullets.append(("section", strip_md_inline(line)))
    return bullets


def parse_notes(p):
    lines = []
    for line in p.split("\n"):
        if line.startswith("> "):
            lines.append(line[2:].rstrip())
        elif line.strip() == ">":
            lines.append("")
    # Collapse: blank line = paragraph break
    text = "\n".join(lines)
    text = re.sub(r"\n+", " ", text).strip()
    # Re-introduce paragraph breaks from explicit blank-line markers
    paragraphs = []
    cur = []
    for line in lines:
        if line.strip() == "":
            if cur:
                paragraphs.append(" ".join(cur))
                cur = []
        else:
            cur.append(line.strip())
    if cur:
        paragraphs.append(" ".join(cur))
    return "\n\n".join(paragraphs)


def strip_md_inline(text):
    """Strip markdown inline formatting (keep text)."""
    # links [text](url) → text
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    # bold **text** → text (kept visually plain; pptx body is uniform)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    # italic *text* → text
    text = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", text)
    # inline code `text` → text
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text


# ── PPTX building ────────────────────────────────────────
def resolve_image(path):
    """Convert a markdown image src to a filesystem path; SVG → PNG."""
    if path.startswith("/"):
        fs = REPO + path
    else:
        fs = os.path.join(os.path.dirname(POST), path)
    if fs.endswith(".svg"):
        return svg_to_png(fs)
    return fs


def add_title(slide, title, top=MARGIN_TOP):
    box = slide.shapes.add_textbox(
        MARGIN_L, top, SLIDE_W - MARGIN_L - MARGIN_R, TITLE_H)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(TITLE_PT)
    run.font.bold = True
    run.font.color.rgb = INK


def add_bullets(slide, bullets, top, height):
    box = slide.shapes.add_textbox(
        MARGIN_L, top, SLIDE_W - MARGIN_L - MARGIN_R, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for depth, text in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if depth == "section":
            # Section header (bold, larger)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(BODY_PT + 2)
            run.font.bold = True
            run.font.color.rgb = INK
        else:
            p.level = depth
            run = p.add_run()
            run.text = "• " + text
            run.font.size = Pt(BODY_PT - depth * 2)
            run.font.color.rgb = INK


def add_images(slide, image_paths, top, max_height, max_width):
    """Place 1-2 images, scaled to fit within (max_width, max_height)."""
    if not image_paths:
        return
    max_h_emu = int(max_height) if isinstance(max_height, int) else max_height.emu
    if len(image_paths) == 1:
        place_image(slide, image_paths[0],
                    Inches(0), top, SLIDE_W, Emu(max_h_emu), center_x=True)
    else:
        # Stack vertically with a small gap
        gap = Inches(0.1)
        each_h = Emu((max_h_emu - gap.emu) // 2)
        top_emu = int(top) if isinstance(top, int) else top.emu
        for i, p in enumerate(image_paths):
            place_image(slide, p, Inches(0),
                        Emu(top_emu + i * (each_h.emu + gap.emu)),
                        SLIDE_W, each_h, center_x=True)


def place_image(slide, png_path, x_unused, y, max_w, max_h, center_x=False):
    """Add an image scaled to fit max_w × max_h, optionally centered horizontally."""
    from PIL import Image
    with Image.open(png_path) as im:
        iw, ih = im.size
    max_w_emu = int(max_w) if isinstance(max_w, int) else max_w.emu
    max_h_emu = int(max_h) if isinstance(max_h, int) else max_h.emu
    y_emu = int(y) if isinstance(y, int) else y.emu
    ratio = min(max_w_emu / iw, max_h_emu / ih) * 0.95  # tiny padding
    w = Emu(int(iw * ratio))
    h = Emu(int(ih * ratio))
    if center_x:
        x = Emu((SLIDE_W.emu - w.emu) // 2)
    else:
        x = MARGIN_L
    y_centered = Emu(y_emu + (max_h_emu - h.emu) // 2)
    slide.shapes.add_picture(png_path, x, y_centered, width=w, height=h)


def add_notes(slide, notes_text):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(NOTE_PT)


def build_title_slide(prs, sl):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Big title (centered)
    h1_h = Inches(1.5)
    box = slide.shapes.add_textbox(
        MARGIN_L, Inches(2.0), SLIDE_W - MARGIN_L - MARGIN_R, h1_h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = sl["title"]
    run.font.size = Pt(44); run.font.bold = True; run.font.color.rgb = INK

    # Subtitle
    box2 = slide.shapes.add_textbox(
        MARGIN_L, Inches(3.6), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.6))
    tf2 = box2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sl["subtitle"]
    r2.font.size = Pt(24); r2.font.italic = True
    r2.font.color.rgb = RGBColor(0x73, 0x7A, 0x82)

    # Byline (author + venue + date)
    box3 = slide.shapes.add_textbox(
        MARGIN_L, Inches(5.5), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.5))
    tf3 = box3.text_frame; tf3.word_wrap = True
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = sl["byline"]
    r3.font.size = Pt(18); r3.font.color.rgb = INK

    add_notes(slide, sl["notes"])


def build_content_slide(prs, sl):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    add_title(slide, sl["title"])

    body_top = Inches(1.2)
    body_h = Emu(SLIDE_H - body_top - Inches(0.3))

    images = [resolve_image(p) for p in sl["images"]]
    bullets = sl["bullets"]

    if images and bullets:
        # Image gets ~65% of body, bullets ~35%
        img_h = Emu(int(body_h.emu * 0.65))
        bul_top = Emu(body_top.emu + img_h.emu + Inches(0.1).emu)
        bul_h = Emu(body_h.emu - img_h.emu - Inches(0.1).emu)
        add_images(slide, images, body_top, img_h, SLIDE_W)
        add_bullets(slide, bullets, bul_top, bul_h)
    elif images:
        add_images(slide, images, body_top, body_h, SLIDE_W)
    elif bullets:
        add_bullets(slide, bullets, body_top, body_h)

    add_notes(slide, sl["notes"])


def main():
    slides = parse_post(POST)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for sl in slides:
        if sl.get("is_title"):
            build_title_slide(prs, sl)
        else:
            build_content_slide(prs, sl)
        print(f"  slide {sl['num']}: {sl['title'][:60]}")

    prs.save(OUT)
    print(f"\nWrote {OUT}")


if __name__ == "__main__":
    main()
